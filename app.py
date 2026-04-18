from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
from groq import Groq
import os
from werkzeug.utils import secure_filename
import pypdf
from pptx import Presentation
import json
from datetime import datetime

app = Flask(__name__, template_folder='templates')
CORS(app)

# ====================== Groq API Key ======================
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")
client = Groq(api_key=GROQ_API_KEY)

LOCAL_KB_PATH = os.path.join(os.path.dirname(__file__), "knowledge_base")
os.makedirs(LOCAL_KB_PATH, exist_ok=True)


class SemanticNetwork:
    def __init__(self):
        self.concepts = {}
        self._initialize_core_knowledge()

    def _initialize_core_knowledge(self):
        """
        Comprehensive APU Knowledge Base.
        Covers: ICS (Integrated Computer Systems), ISFT (Info Security & Forensic Tech),
                SAAD (Systems Analysis and Design), IAI (Intro to Artificial Intelligence)
        Format: (concept_key, subject_chapter_topic, detailed_description)
        """
        core_concepts = [
            ("ICS_CH1_Integrated_Systems_Definition", "ICS | CH1 - Integrated Systems Definition", "Integrated systems are combinations of hardware, software, and peripheral components that work together cohesively to perform complex tasks. Designed to operate as a single unit with seamless communication. Complex tasks require multiple steps, involve large data processing, need coordination between hardware/software/networks, and demand automation or parallel processing."),
            ("ICS_CH1_Efficiency", "ICS | CH1 - Significance: Efficiency", "Efficiency in ICS means achieving maximum performance while using least possible resources (time, energy, cost). Includes: Processing Efficiency (fast CPU/GPU task completion, high throughput), Resource Utilization (optimal use of memory/storage/bandwidth, avoiding bottlenecks), Energy Efficiency (minimal power, critical for mobile/IoT/data centers), User-Level Efficiency (smooth applications, minimal delays). Achieving maximum performance while minimizing resource usage."),
            ("ICS_CH1_Scalability", "ICS | CH1 - Significance: Scalability", "Scalability is the system ability to grow and handle increasing workloads without losing performance. Vertical scaling = upgrading hardware (more RAM, faster CPU). Horizontal scaling = adding more systems/servers. Properties: Flexibility (adapt without full redesign), Performance Maintenance (runs smoothly even when tasks double), Cost-Effectiveness (add resources only when needed). Example: Lazada/Shopee adds servers during festive sales to handle millions of customers."),
            ("ICS_CH1_Innovation", "ICS | CH1 - Significance: Innovation", "Innovation in ICS is creating new technologies by integrating computing components into unified systems. Drives advancements by combining AI, IoT, cloud, big data. Enables new applications: remote patient monitoring, AI-assisted diagnosis, telemedicine. Improves existing systems: banking ICS with mobile apps, biometric security, real-time fraud detection. Cross-industry impact: education, manufacturing, transportation. Example: autonomous cars use integrated sensors, AI, and cloud connectivity."),
            ("ICS_CH1_Reliability", "ICS | CH1 - Significance: Reliability", "Reliability in ICS ensures the system performs accurately, consistently, and without failure. Properties: Consistency (same behavior for same inputs), Fault Tolerance (keeps working even if component fails), Availability (accessible when users need it), Error Handling (detects and recovers from errors gracefully), Durability (data not lost/corrupted). Examples: Banking systems process 24/7 without downtime; Hospital systems monitor patients consistently; Google Drive ensures files accessible even if a server fails."),
            ("ICS_CH1_System_Architecture", "ICS | CH1 - System Architecture and Advanced Systems Software", "System architecture is the structured framework defining components, relationships, and interactions within an integrated system. Includes Hardware Components (processors, sensors, controllers), Software Components (programs controlling hardware), Communication Interfaces (protocols enabling data exchange). Advanced Systems Software manages and optimizes hardware/software resources. Key types: Operating Systems (Windows, macOS, Linux, Unix), Device Drivers (NVIDIA, Realtek, Printer drivers), System Utilities (Antivirus like Norton, Disk Cleanup like CCleaner), Firmware (BIOS, UEFI), Shells and CLI (Bash, PowerShell), Hypervisors (VMware ESXi, Microsoft Hyper-V, Oracle VirtualBox)."),
            ("ICS_CH2_Computer_Integration", "ICS | CH2 - Computer Integration Definition and Types", "Computer Integration (CI) is the process of combining various computing systems and software to operate as a cohesive unit. Three types: 1) Vertical Integration - combines systems within a specific domain/industry. Example: EHR (Electronic Health Records) integrates hospital, clinic, and lab patient data into one system. 2) Horizontal Integration - connects systems across different domains. Example: ERP connects finance, HR, supply chain, CRM across departments. Smart Cities integrate traffic, transport, energy, waste management. 3) Data Integration - combines data from different sources for seamless consistent flow. Example: CRM integrating customer data from online stores, physical outlets, and social media."),
            ("ICS_CH2_CI_Benefits_Challenges", "ICS | CH2 - Computer Integration Benefits, Components, and Challenges", "Benefits: Operational Efficiency (streamlines processes, reduces redundancy), Cost-Effectiveness (optimizes resource utilization), Improved Communication (better data exchange and collaboration). Components: Hardware Integration (combining servers, networks, storage), Software Integration (merging applications - ERP integrating finance/HR/supply chain), Network Integration (connecting networks via VPNs, smart home hubs). Challenges: Compatibility Issues (legacy systems need middleware), Network Security Concerns (different encryption standards), Human Factors (resistance to change, need for training). Future Trends: AI and Machine Learning, Internet of Things (IoT), Cloud Computing."),
            ("ICS_CH3_Von_Neumann_Architecture", "ICS | CH3 - Von Neumann Architecture", "Von Neumann Architecture (1945-1951) introduced the stored-program concept: both data and instructions stored together in same memory. Before this, computers were hardwired (ENIAC needed rewiring 18000 vacuum tubes for each new task). Guidelines: Stored Program Concept (same memory for instructions and data), Linear Addressing (sequential unique address for each memory location), Independence of Data and Address (address tells where, not what), Sequential Processing (instructions processed in sequence using Program Counter), Functional Organization (CPU = Control Unit + ALU + Registers). Components: Memory, CPU, I/O devices. Fetch-Decode-Execute cycle: Fetch instruction from memory, Decode it, Execute in ALU, Store results. Still the dominant architecture for modern computers."),
            ("ICS_CH3_Harvard_Architecture", "ICS | CH3 - Harvard Architecture vs Von Neumann", "Harvard Architecture has separate memory and buses for instructions and data, allowing simultaneous access. Two buses: one for instructions, one for data, preventing bottlenecks. Used in microcontrollers and DSPs (Digital Signal Processors). Advantages: Increased Speed (simultaneous instruction fetch and data access), Better Performance for DSP and embedded systems, Security against buffer overflow attacks. Modern CPUs use modified Harvard architecture internally (separate L1 caches) while appearing as Von Neumann externally. Comparison: Von Neumann = single memory and bus, simpler/cheaper, used in PCs/laptops/servers, sequential fetch. Harvard = separate memories and buses, faster/more complex/expensive, used in embedded systems/DSPs/microcontrollers, simultaneous fetch."),
            ("ICS_CH4_Instruction_Set_CISC_RISC", "ICS | CH4 - Instruction Set Architecture, CISC and RISC", "Instruction Set is the vocabulary of the processor - all instructions it understands. Elements of an instruction: OP-CODE (what task to do), Source OPERAND(s) (input data), Result OPERAND (where to store output). Example: ADD R1, R2, R3 - OP-CODE=ADD, Source=R1 and R2, Result=R3. Addressing modes: Direct (instruction points to memory address), Indirect (points to address which has the real address), Register (data in register), Immediate (data in instruction itself). CISC (Complex Instruction Set Computer): reduces instructions per program, complex hardware, simpler software, variable instruction size, microcoded control, fewer registers. Example: Intel x86 (PCs, laptops, servers). RISC (Reduced Instruction Set Computer): reduces cycles per instruction, simpler hardware, complex compiler, fixed instruction size, hardwired control, large number of registers, pipelining support. Example: ARM (smartphones, Apple M1/M2/M3, IoT). CISC = Swiss Army knife; RISC = set of specialized simple tools."),
            ("ICS_CH4_Pipelining_and_Hazards", "ICS | CH4 - Pipelining, Data Dependency and Branch Prediction", "Pipelining overlaps instruction execution to reduce overall execution time. Stages: Fetch Instruction (FI), Decode Address (DA), Fetch Operand (FO), Execute (EX). Like an assembly line - while one instruction executes, next is decoded, next after is fetched. Problems: 1) Data Dependency - two adjacent instructions access same memory location simultaneously. Solutions: Hardware Interlocks (CPU makes next instruction wait like a traffic light), Operand Forwarding/Bypassing (CPU directly passes result to next instruction without waiting), Delayed Load (compiler rearranges instructions to fill wait time). 2) Branching of Instructions - when branch occurs, pre-fetched instructions in pipeline must be discarded (pipeline stall, wrong instructions loaded). Solutions: Stalling (wait until branch decision known - wastes cycles), Branch Prediction (CPU guesses path - static: fixed rule; dynamic: learns from history), Branch Delay Slots (instruction after branch always executes, compiler fills it usefully), Speculative Execution (executes along predicted path, discards if wrong - security risk: Spectre, Meltdown)."),
            ("ICS_CH4_CPU_Registers_Flags", "ICS | CH4 - CPU Registers and Status Flags (PSW)", "Registers are small fast storage locations inside the CPU. Types: General-Purpose Registers (eax, ebx, ecx, edx in x86 - used for various operations; ax=lower 16 bits of eax; ah=upper 8 bits of ax; al=lower 8 bits of ax), Special-Purpose Registers (MAR - Memory Address Register stores memory address to read/write; MDR - Memory Data Register stores data being transferred; Program Counter/PC - address of next instruction; Instruction Register - current instruction being executed). Status Flags (1-bit cells in Program Status Word/PSW): Carry (CY) Flag - set if carry from MSB in arithmetic operation; Parity (P) Flag - set if result has even number of 1 bits; Auxiliary-Carry (AC/HC) Flag - set if carry from bit 3 to bit 4, used in BCD arithmetic; Zero (Z) Flag - set if result of arithmetic operation is zero; Sign (S) Flag - set if result is negative (MSB=1); Overflow (V) Flag - set if signed overflow occurs, result too large for available bits."),
            ("ICS_CH5_Memory_Hierarchy", "ICS | CH5 - Memory Hierarchy and Types", "Memory hierarchy from fastest/smallest to slowest/largest: 1) Registers (inside CPU, fastest, few bytes, stores immediate data), 2) Cache Memory (L1 fastest inside CPU core, L2 slightly larger/slower, L3 shared among cores, stores frequently used data), 3) Main Memory RAM (stores programs and data in use, volatile), 4) Local Secondary Storage (HDD, SSD, Flash Drives, non-volatile, large capacity), 5) Remote/Cloud Storage (backup, archival, slowest, very large). Two divisions: Internal/Primary Memory (Main Memory + Cache + Registers, CPU directly accessible) and External/Secondary Memory (Magnetic Disk, Optical Disk, Magnetic Tape, accessed via I/O Module). RAM types: SRAM (Static RAM, faster/expensive, no refresh needed, used in cache), DRAM (Dynamic RAM, slower/cheaper, needs refresh every ~4ms, used in main memory). DDR SDRAM transfers on both clock edges (double pumping). DDR2 performs 4 transfers per clock using multiplexing on wider bus. ROM types: PROM (program once), EPROM (erase with UV light), EEPROM (erase electrically), Flash ROM (reprogram while in computer, used in USB/SSD)."),
            ("ICS_CH5_Endianness_Memory_Capacity", "ICS | CH5 - Endianness and Memory Capacity", "Endianness describes byte storage order in memory. Big-Endian: most significant byte (MSB) stored first at lowest address. Little-Endian: least significant byte (LSB) stored first. Example: 32-bit value 0x12345678 - Big-Endian stores 12 34 56 78; Little-Endian stores 78 56 34 12. Intel x86 uses Little-Endian. Network byte order is Big-Endian. Important for cross-platform data exchange and assembly programming. Memory capacity units: bit (smallest, 0 or 1), nibble = 4 bits, byte = 8 bits, kilobyte (KB) = 1024 bytes, megabyte (MB) = 1024 KB, gigabyte (GB) = 1024 MB, terabyte (TB) = 1024 GB, petabyte (PB) = 1024 TB. Storage types: Magnetic (HDD, Floppy, Magnetic Tape), Optical (CD/DVD), Flash Memory (SSD, USB drives), Online Cloud Storage (Google Drive, AWS S3)."),
            ("ICS_CH6_Deadlock", "ICS | CH6 - Deadlock in ICS", "Deadlock occurs when threads wait indefinitely for resources held by each other, making no progress. Classic example: Dining Philosophers Problem. Four necessary conditions (ALL must exist simultaneously): 1) Mutual Exclusion (at least one resource is non-sharable), 2) Hold and Wait (thread holding resources waits for more held by others), 3) No Preemption (resources cannot be forcibly taken, must be released voluntarily), 4) Circular Wait (T1 waits for T2, T2 waits for T3, ..., Tn waits for T1). Resource-Allocation Graph: threads=circles, resources=rectangles with dots. Request edge Ti→Rj, Assignment edge Rj→Ti. No cycle = no deadlock. Cycle with one instance per resource = definite deadlock. Cycle with multiple instances = possible deadlock. Three handling methods: 1) Ignore (most OS use this, cheaper), 2) Prevent or Avoid (Deadlock Prevention: ensures at least one condition cannot hold; Deadlock Avoidance: uses advance info, Banker Algorithm checks for safe state), 3) Detect and Recover (detect via wait-for graph, recover by aborting threads or preempting resources)."),
            ("ICS_CH7_Processes_Threads", "ICS | CH7 - Processes, Threads, and Multithreading", "Process: a program in execution with its own isolated memory space, managed by OS, slower to create. Lifecycle states: New, Ready, Running, Waiting, Terminated. IPC (Inter-Process Communication): Shared Memory (fast, direct memory access, needs semaphore/mutex synchronization, local only) vs Message Passing (through OS kernel, safer, local or networked, OS handles sync). Thread: smallest unit of execution within a process, shares process memory, has own stack and registers, faster to create than processes. Process vs Thread: Processes have separate memory (isolated, safer), threads share memory (faster communication, risk of race conditions). Multithreading models: 1) Many-to-One (many user threads → one kernel thread, if one blocks all block, not parallel), 2) One-to-One (each user thread → one kernel thread, true parallelism, high overhead, used in Windows/Linux), 3) Many-to-Many (user threads mapped to fewer kernel threads dynamically, best balance, used in Solaris 9). Synchronization tools: Mutexes, Semaphores, Monitors, Locks prevent race conditions."),
            ("ISFT_CH1_Digital_Footprint_Terms", "ISFT | CH1 - Digital Footprint and Security Terms", "Digital Footprint: data trail you leave online. Active footprint: deliberate actions (posting on social media, submitting forms, online banking). Passive footprint: data collected without direct knowledge (cookies tracking visits, apps collecting location, social media behavior tracking). Key security terms: Vulnerability (weakness that can be exploited), Threat (potential danger to system), Attack (intentional action exploiting vulnerability), Exploit (method/tool to take advantage of vulnerability), Risk (probability that threat exploits vulnerability and causes harm - Risk = Threat x Vulnerability x Impact), Incident (confirmed security event violating security policy), Incident Response (process: identify, contain, eradicate, recover, learn)."),
            ("ISFT_CH1_CIA_Triad_AAA", "ISFT | CH1 - CIA Triad and AAA Framework", "CIA Triad - three core information security principles: 1) Confidentiality - information accessible only to authorized users. Methods: encryption, access controls, authentication, data classification. Threats: unauthorized disclosure, eavesdropping. 2) Integrity - information is accurate and unaltered. Methods: hashing (MD5, SHA-256), digital signatures, checksums. Threats: data tampering, man-in-the-middle attacks. 3) Availability - systems and information accessible when needed. Methods: redundancy, backups, failover, DDoS protection, load balancing. Threats: DDoS attacks, ransomware, hardware failures. AAA Framework: Authentication (verify identity - something you know/have/are, MFA combines 2+ factors), Authorization (determine what user can do - ACL, RBAC, MAC), Accounting/Auditing (track what authenticated users do - log files, audit trails, SIEM)."),
            ("ISFT_CH2_Defense_in_Depth", "ISFT | CH2 - Layered Security and Defense in Depth", "Defense in Depth (DiD): multiple layers of security controls so if one fails others still protect. Layers from outside to inside: Physical Security (locks, cameras, guards, biometric access), Network Security (firewalls, IDS/IPS, VPNs, network segmentation), Host Security (antivirus, host-based IDS, OS hardening, patch management), Application Security (secure coding, WAF, input validation, authentication), Data Security (encryption, DLP, backups, access controls). Security model flow: Prevention → Detection → Response → Recovery. No single control is sufficient. Principle of least privilege: give users minimum access needed. Need to know principle: only access data relevant to role. Security through obscurity alone is not sufficient."),
            ("ISFT_CH3_Cybercrime_Dark_Web", "ISFT | CH3 - Cybercrime and the Dark Web", "Cybercrime categories: Against individuals (identity theft, cyberstalking, phishing, fraud), Against organizations (data breaches, ransomware, corporate espionage), Against governments (cyberterrorism, infrastructure attacks). Dark Web: part of internet not indexed by search engines, accessible only through Tor Browser. Used for: anonymous communication, buying stolen credentials/illegal goods, ransomware-as-a-service, hiring hackers. Tor Network: routes traffic through multiple encrypted relays (nodes) to anonymize users. Three levels of internet: Clearnet/Surface Web (normal internet), Deep Web (requires login - email, banking), Dark Web (requires Tor/special software). Cybercrime investigation requires: evidence preservation, chain of custody, legal authorization, following proper forensic methodology."),
            ("ISFT_CH4_IAM", "ISFT | CH4 - Identity and Access Management (IAM)", "IAM is a framework ensuring right users have appropriate access to technology resources. Components: Identity Management (creating, managing, deleting user identities - provisioning/de-provisioning), Access Management (controlling resource access), Authentication (verifying identity), Authorization (what you can do), Single Sign-On (SSO - one login for multiple systems), Privileged Access Management (PAM - controls/monitors admin access). Common attacks: Credential stuffing (using leaked username/password lists from data breaches), Brute force (trying all combinations), Pass-the-hash (using captured password hash without knowing plaintext), Privilege escalation (gaining higher permissions than authorized). Implementations: LDAP (Lightweight Directory Access Protocol), Microsoft Active Directory. MFA factors: Something you know (password/PIN), Something you have (smart card/token/OTP), Something you are (biometrics)."),
            ("ISFT_CH5_Cryptography", "ISFT | CH5 - Cryptography and PKI", "Cryptography secures information by transforming it to unreadable format. Key terms: Plaintext (original data), Ciphertext (encrypted data), Key (secret value used for encryption/decryption). Symmetric encryption: same key for encrypt/decrypt, fast, key distribution problem. Examples: AES (most secure, 128/192/256-bit), DES (outdated), 3DES. Asymmetric encryption: public key encrypts, private key decrypts, slower but solves key distribution. Examples: RSA, ECC, Diffie-Hellman. Hashing: one-way function producing fixed-length digest for integrity verification. MD5 (128-bit, broken), SHA-1 (160-bit, deprecated), SHA-256/SHA-3 (secure). Digital Signature: hash encrypted with private key - proves authenticity, integrity, non-repudiation. PKI (Public Key Infrastructure): Certificate Authority (CA) issues/signs digital certificates. X.509 certificate contains: public key, owner identity, CA signature, validity period. SSL/TLS uses PKI for HTTPS. Certificate Revocation: CRL (Certificate Revocation List), OCSP (Online Certificate Status Protocol)."),
            ("ISFT_CH6_Policies_Compliance", "ISFT | CH6 - Security Policies, Standards and Compliance", "Security Policy: formal document defining security rules and responsibilities. Types: Acceptable Use Policy (AUP), Password Policy, Data Classification Policy, Incident Response Policy, BYOD Policy. Standards and Frameworks: ISO/IEC 27001 (Information Security Management System - ISMS), NIST Cybersecurity Framework (Identify, Protect, Detect, Respond, Recover), PCI-DSS (Payment Card Industry - for organizations handling credit cards), HIPAA (Health Insurance Portability and Accountability Act - US healthcare data), GDPR (General Data Protection Regulation - EU data protection, fines up to 4% of global revenue). Data Classification levels: Public (anyone can see), Internal (employees only), Confidential (need-to-know), Secret/Top Secret (highest sensitivity). Non-compliance consequences: fines, legal action, reputational damage."),
            ("ISFT_CH7_Risk_Management", "ISFT | CH7 - Risk Management and Privacy", "Risk Management process: 1) Risk Identification (identify assets, threats, vulnerabilities), 2) Risk Assessment (calculate likelihood and impact), 3) Risk Treatment - 4 options: Accept (tolerate risk), Avoid (eliminate activity causing risk), Transfer (insurance, outsourcing), Mitigate (reduce likelihood/impact), 4) Risk Monitoring (continuously review/update). Quantitative Analysis: ALE = ARO x SLE (Annual Loss Expectancy = Annual Rate of Occurrence x Single Loss Expectancy). Qualitative Analysis: High/Medium/Low scales. Privacy principles (GDPR): Data minimization (collect only what needed), Purpose limitation (use only for stated purpose), Storage limitation (don't keep longer than needed), Transparency (inform users). GDPR rights: Right to access, Right to rectification, Right to erasure/be forgotten, Right to data portability."),
            ("ISFT_CH8_Malware_Attacks", "ISFT | CH8 - Malware Types and Common Attack Methods", "Malware types: Virus (attaches to programs, spreads when program runs, needs human action), Worm (self-replicating, spreads through networks automatically), Trojan Horse (disguised as legitimate software, creates backdoor), Ransomware (encrypts files, demands ransom - WannaCry, CryptoLocker), Spyware (secretly monitors activity), Adware (unwanted ads), Rootkit (hides at OS/kernel level, very hard to detect), Keylogger (records keystrokes to capture passwords), Botnet (network of infected computers for DDoS/spam), RAT (Remote Access Trojan - gives full remote control). Attack types: Phishing (fraudulent emails/websites stealing credentials), Spear Phishing (targeted), Whaling (targeting executives), Vishing (voice phishing). SQL Injection (malicious SQL in input fields: ' OR '1'='1, prevention: parameterized queries). XSS/Cross-Site Scripting (injecting JavaScript into pages, types: Stored/Reflected/DOM-based). Man-in-the-Middle (intercepting communication). DDoS (overwhelming server with botnet traffic). Social Engineering (manipulation: Pretexting, Baiting, Tailgating)."),
            ("ISFT_CH9_Digital_Forensics", "ISFT | CH9 - Digital Forensics Methodology", "Digital Forensics: collecting, preserving, analyzing, and presenting digital evidence in legally admissible manner. 5 Phases: 1) Identification (identify evidence sources: computers, mobile devices, networks, cloud), 2) Preservation (secure evidence, create forensic image - bit-by-bit copy, use write blockers to prevent changes, calculate hash MD5/SHA-256 before and after), 3) Collection (gather evidence following proper procedures, maintain Chain of Custody - documentation of who had access to evidence and when), 4) Analysis (examine evidence using tools: Autopsy, Volatility for RAM, Wireshark for network, FTK Forensic Toolkit), 5) Reporting (present findings clearly and defensibly). Order of Volatility (most volatile first): CPU registers/cache, RAM, Network connections, Running processes, Disk, Optical media. RAM Imaging tools: RAM Magnet, DumpIt, WinPmem. Volatility Framework commands: imageinfo, pslist (list processes), netscan (network connections), filescan (files in memory), malfind (find injected code). HDD Imaging: FTK Imager, dd command. Image formats: E01 (EnCase), dd/raw, AFF. Autopsy: analyzes disk images, recovers deleted files, timeline analysis, keyword search."),
            ("ISFT_CH10_Network_Forensics_Steganography", "ISFT | CH10 - Network Forensics, Wireshark and Steganography", "Wireshark: most popular network protocol analyzer/packet sniffer. Captures packets in real-time, saves in .pcap format. Key filters: ip.addr==x.x.x.x (filter by IP), tcp.port==80 (filter by port), http (HTTP traffic), dns (DNS queries), follow TCP stream (reconstruct conversation). Network forensics indicators: unusual port numbers, large data transfers, connections to known malicious IPs, DNS queries to suspicious domains, cleartext passwords (HTTP/FTP/Telnet). Network attacks visible in Wireshark: ARP poisoning (duplicate ARP replies), port scanning (many SYN packets to different ports), DDoS (traffic flood), DNS tunneling (unusually large DNS packets). Windows Registry forensics: HKLM\\SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\Run (startup programs), HKCU\\Software\\Microsoft\\Windows\\CurrentVersion\\Explorer\\RecentDocs (recent files), registry reveals malware persistence. File Signatures (Magic Bytes): PDF=25504446, JPEG=FFD8FF, PNG=89504E47, ZIP=504B0304. Steganography: hiding data inside ordinary files (images, audio). LSB method modifies least significant bits of pixels. Detection (Steganalysis): statistical analysis, compare file size. Tools: Steghide, OpenStego."),
            ("SAAD_CH1_Information_Systems", "SAAD | CH1 - Information Systems Definition and Types", "Information System (IS): set of interrelated components that collect, process, store, and distribute information to support decision making. 5 Components: Hardware (servers, workstations, networks, scanners, mobile devices), Software (programs controlling hardware), Data (raw facts transformed into useful information, data alone has no meaning), Processes (tasks users/managers/IT perform to achieve results), People (users and stakeholders who communicate with IS). Types of IS: Transaction Processing Systems (TPS - day-to-day routine transactions, mission-critical: payroll, billing, stock control), Management Information Systems (MIS - converts TPS data into reports for middle managers), Decision Support Systems (DSS - analytical models helping semi-structured decisions), Executive Support Systems (ESS - strategic info for senior executives), Enterprise Resource Planning (ERP - integrates all departments: SAP), Knowledge Management Systems/Expert Systems (simulate human reasoning using knowledge base and inference rules, fuzzy logic), User Productivity Systems (groupware: Calendar, Email, Scheduling)."),
            ("SAAD_CH2_SDLC_Methodologies", "SAAD | CH2 - SDLC and Development Methodologies", "SDLC (Systems Development Life Cycle) has 5 phases: 1) Systems Planning (identify and prioritize projects, feasibility study), 2) Systems Analysis (study existing system, gather requirements, create AS-IS models, deliverable: System Requirements Specification), 3) Systems Design (design new system, create TO-BE models, database/UI/architecture design, deliverable: Design Specification), 4) Systems Implementation (build/code, test, deploy), 5) Systems Support and Security (maintain, update, secure). Development Methodologies: Structured Analysis (traditional sequential SDLC, uses DFDs and ERDs), Object-Oriented Analysis (data + processes = objects, object has class, properties/attributes, methods/behaviors, objects communicate via messages), Agile Methods (continuous feedback, iterative development, Agile Manifesto: individuals over processes, working software over documentation, customer collaboration over contract, responding to change over following plan). Agile frameworks: Scrum (sprints 2-4 weeks, daily standups, Product Backlog, Sprint Backlog, Retrospectives), XP/Extreme Programming (pair programming, test-driven development). Spiral Model (iterative with risk assessment). JAD (Joint Application Development - workshops with stakeholders). RAD (Rapid Application Development - prototyping + iterative delivery)."),
            ("SAAD_CH2_Feasibility_Analysis", "SAAD | CH2 - Feasibility Analysis", "Feasibility Study determines if proposed system should be built. Four types: 1) Technical Feasibility - does technology exist? Can we build it with available hardware/software? Is team capable? 2) Economic Feasibility/Cost-Benefit Analysis - do benefits outweigh costs? Tangible benefits (measurable: increased sales, reduced costs), Intangible benefits (hard to measure: improved morale, better customer service). Calculate ROI (Return on Investment) and payback period. 3) Operational Feasibility - will users accept the system? Are business processes compatible? Is management supportive? 4) Schedule Feasibility - can system be built within required timeframe? Available resources? Sufficient skilled staff?"),
            ("SAAD_CH3_Requirements_Gathering", "SAAD | CH3 - Requirements Gathering Techniques", "Methods for gathering system requirements: 1) Interviews (direct conversations with stakeholders, types: structured/unstructured/semi-structured, best for detailed information but time-consuming), 2) Questionnaires/Surveys (written questions for large groups, statistical analysis possible, lower response rate), 3) Observation (watching users in actual environment, reveals unstated/actual requirements), 4) Document Review (analyzing existing reports, forms, manuals, reveals current data flows), 5) JAD Sessions (workshops with all stakeholders together), 6) Prototyping (building preliminary version - throwaway prototype for clarification only, evolutionary prototype becomes final system). Use Case Analysis: identifying system actors and their interactions."),
            ("SAAD_CH3_DFD_ERD", "SAAD | CH3 - Data Flow Diagrams (DFD) and Entity-Relationship Diagrams (ERD)", "DFD components: Process (circle/bubble - transforms data, labeled with verb phrase), External Entity (rectangle - data source/destination outside system boundary e.g. Customer, Manager), Data Store (parallel lines - where data stored e.g. D1 Student Records), Data Flow (arrow with label - data moving between components). DFD Levels: Context Diagram (Level 0 - single process for entire system, all external entities and flows), Level 1 (decomposes context diagram into major processes), Level 2 (further decomposes Level 1 processes). Rules: flows to/from processes only, no unnamed flows, balanced decomposition. ERD components: Entity (rectangle - thing/object e.g. Student, Course), Attribute (oval - property of entity e.g. StudentID, Name), Relationship (diamond - association e.g. Enrolls), Primary Key (underlined attribute - unique identifier). Cardinality: One-to-One (1:1), One-to-Many (1:M), Many-to-Many (M:N). Crow's Foot Notation: single line (one), crow's foot (many), circle (zero/optional), vertical line (mandatory). Normalization: 1NF (no repeating groups, atomic values), 2NF (no partial dependencies), 3NF (no transitive dependencies)."),
            ("SAAD_CH4_Use_Case_UML", "SAAD | CH4 - Use Case Diagrams and UML", "Use Case Diagram (UML) shows functional requirements from user perspective. Components: Actor (stick figure - person/system interacting with system), Use Case (oval - function/service system provides), System Boundary (rectangle - defines inside/outside system), Relationships: Association (solid line actor to use case), Include (dashed arrow with include - mandatory shared behavior), Extend (dashed arrow with extend - optional behavior under certain conditions), Generalization (inheritance between actors or use cases). Other UML diagrams: Class Diagram (system structure - classes, attributes, methods, relationships), Sequence Diagram (object interactions over time - who sends what message to whom and when), Activity Diagram (flowchart showing workflow/business process logic). CRUD matrix: shows which processes Create, Read, Update, Delete each data entity."),
            ("SAAD_CH5_Systems_Design", "SAAD | CH5 - Systems Design and Interface Design", "Input Design principles: minimize keystrokes, validate data at entry, use defaults, provide feedback. Output Design: clarity, accuracy, completeness, timeliness. Types: Scheduled reports (on schedule), Exception reports (unusual conditions), Demand reports (on request). Database Design: convert ERD to relational tables, define primary/foreign keys, ensure normalization, define data types/constraints. Interface Design (Nielsen's 10 Heuristics): Visibility of system status, Match system to real world, User control and freedom, Consistency and standards, Error prevention, Recognition over recall, Flexibility and efficiency, Aesthetic and minimalist design, Help users recognize/diagnose/recover from errors, Help and documentation. UI types: CLI (Command Line Interface), GUI (Graphical User Interface), NUI (Natural User Interface - touch, voice, gesture). Prototyping approaches: paper prototypes, wireframes (low-fidelity), mockups (high-fidelity)."),
            ("SAAD_CH6_Project_Management", "SAAD | CH6 - Project Management for IS Development", "Project Management: planning, organizing, controlling resources to achieve goals within constraints. Triple Constraint/Project Triangle: Scope, Time (Schedule), Cost - changing one affects others. Quality sometimes added as 4th constraint. Project Management tools: Gantt Chart (bar chart showing tasks against timeline, task dependencies, milestones, deadlines), PERT Chart (network diagram showing task dependencies and critical path), Critical Path Method (CPM - identifies longest sequence of dependent tasks determining project duration). WBS (Work Breakdown Structure): hierarchical decomposition of project into manageable components. Risk Management in projects: identify risks, assess probability/impact, create risk response plan (avoid, mitigate, transfer, accept). Agile PM: sprints, daily standups, velocity tracking, burndown charts showing remaining work vs time. Project Manager responsibilities: scope management, time management, cost management, quality management, risk management, stakeholder management."),
            ("AI_CH1_AI_Definitions_Types", "IAI | CH1 - AI Definitions and Types", "Artificial Intelligence (AI): branch of computer science creating machines that imitate human thinking and reasoning. Definitions: Minsky - 'science of making machines do things requiring intelligence if done by humans'; Feigenbaum - 'part of computer science designing intelligent computer systems'; McCarthy - 'science and engineering of making intelligent machines'; Rouse - 'simulation of human intelligence processes by machines including learning, reasoning, self-correction'. Four goals: Think Humanly (machine behavior matches human), Act Humanly (Turing Test - human-level performance in cognitive tasks), Think Rationally (logical reasoning), Act Rationally (making correct inferences to achieve goals). AI types by capability: Narrow AI/Weak AI (specific task only - Siri, chess programs, spam filters, recommendation systems), General AI/Strong AI (hypothetical human-level AI across all domains), Superintelligence (hypothetical AI surpassing human intelligence). AI types by functionality: Reactive Machines (no memory, react to current input - Deep Blue chess), Limited Memory AI (uses past data - self-driving cars), Theory of Mind AI (hypothetical, understands emotions/beliefs), Self-Aware AI (hypothetical consciousness)."),
            ("AI_CH2_Knowledge_Representation", "IAI | CH2 - Knowledge Representation Methods", "Knowledge Representation (KR): how AI systems store and use knowledge for reasoning. Four methods: 1) Logical Representation (formal logic, Propositional Logic: statements True/False, operators AND/OR/NOT/IMPLIES; Predicate Logic: quantifiers and predicates - All humans are mortal. Advantage: precise and formal. Disadvantage: limited expressiveness). 2) Semantic Networks (graph of nodes/concepts and labeled edges/relationships. Types of relationships: IS-A (inheritance - Dog IS-A Animal), HAS-A (composition - Car HAS-A Engine), RELATED-TO (associations). Advantage: visual, captures relationships. Used in StudBud chatbot). 3) Frames (data structures representing stereotyped situations like object classes. Contains slots/attributes and fillers/values. Inheritance between frames. Example: Car frame has slots: color, brand, speed. Specific car fills slots with values). 4) Production Rules/IF-THEN Rules (condition-action rules. IF condition THEN action. Used in expert systems. Forward Chaining: data-driven, start from facts to find conclusions. Backward Chaining: goal-driven, start from goal and find supporting facts)."),
            ("AI_CH2_Expert_Systems", "IAI | CH2 - Expert Systems", "Expert Systems: AI programs simulating decision-making ability of human experts in specific domain. Components: 1) Knowledge Base (domain-specific facts + rules, curated by knowledge engineers), 2) Inference Engine (applies rules to known facts to deduce new facts, forward or backward chaining), 3) User Interface (allows non-experts to query the system), 4) Explanation Facility (explains how conclusions were reached - crucial for trust), 5) Knowledge Acquisition Facility (tool for adding/updating knowledge). Types of knowledge: Declarative (facts: Python is a programming language), Procedural (how to do things: algorithm steps), Heuristic (rules of thumb: if temperature > 38°C patient likely has fever). Historical expert systems: MYCIN (medical diagnosis of blood infections), XCON (computer configuration). Modern applications: medical diagnosis, financial planning, fault diagnosis, customer support chatbots. Limitations: knowledge acquisition bottleneck, brittle outside domain, cannot learn on own."),
            ("AI_CH3_Search_Algorithms", "IAI | CH3 - Search Algorithms in AI", "Search algorithms find solutions in problem space. Problem formulation: State Space (all possible states), Initial State, Goal State, Actions/Operators, Path Cost. Uninformed (Blind) Search - no knowledge about goal distance: Breadth-First Search (BFS - explores all nodes at current depth before going deeper, complete and optimal for equal costs, O(b^d) time and space, b=branching factor d=depth), Depth-First Search (DFS - explores deep before backtracking, not complete for infinite states, not optimal, O(b*d) space), Depth-Limited Search (DLS - DFS with maximum depth limit), Iterative Deepening DFS (IDDFS - tries increasing depth limits, combines BFS optimality with DFS space efficiency). Informed (Heuristic) Search - uses additional info: Greedy Best-First Search (selects node closest to goal using heuristic h(n), not optimal), A* Search (uses f(n) = g(n) + h(n), g(n)=actual cost from start, h(n)=estimated cost to goal, optimal and complete if h(n) is admissible/never overestimates, most important AI search algorithm). Local Search: Hill Climbing (moves to best neighbor, stuck at local maxima), Simulated Annealing (sometimes accepts worse states to escape local optima), Genetic Algorithms (evolutionary - population, fitness function, selection, crossover, mutation)."),
            ("AI_CH4_Machine_Learning", "IAI | CH4 - Machine Learning and Deep Learning", "Machine Learning (ML): systems learn from data without being explicitly programmed. Types: 1) Supervised Learning (labeled training data, input-output pairs). Classification (spam detection, image recognition). Regression (house prices, stock values). Algorithms: Linear/Logistic Regression, Decision Trees, Random Forest, SVM, k-NN. 2) Unsupervised Learning (unlabeled data, finds patterns). Clustering (customer segmentation - k-Means), Dimensionality Reduction (PCA). 3) Reinforcement Learning (agent learns from environment rewards/penalties, used in game playing - AlphaGo, robotics, autonomous driving). ML Pipeline: Data Collection, Preprocessing (cleaning, normalization, feature engineering), Model Selection, Training, Validation, Testing, Deployment. Overfitting = model too complex, memorizes training data, poor generalization. Fix with regularization, cross-validation, more data. Neural Networks: inspired by biological neurons. Perceptron = single neuron with weights, bias, activation function. ANN has Input Layer, Hidden Layers, Output Layer. Backpropagation trains networks via gradient descent. Activation functions: Sigmoid (0-1, binary classification), ReLU (max(0,x), hidden layers), Softmax (multi-class). Deep Learning = many hidden layers. CNNs for images (convolution + pooling). RNNs for sequences. LSTM solves vanishing gradient. Transformers = basis for LLMs (GPT, BERT, Gemini)."),
            ("AI_CH5_NLP_Computer_Vision", "IAI | CH5 - NLP and Computer Vision", "NLP (Natural Language Processing): AI field enabling machines to understand and generate human language. Key tasks: Tokenization (breaking text into words/sentences), POS Tagging (identifying nouns/verbs/adjectives), Named Entity Recognition/NER (identifying names, places, dates), Sentiment Analysis (positive/negative/neutral), Machine Translation (Google Translate), Text Summarization, Question Answering, Chatbots. NLP Pipeline: Raw Text, Preprocessing (lowercase, remove punctuation/stopwords, stemming/lemmatization), Feature Extraction (Bag of Words, TF-IDF, Word Embeddings like Word2Vec/GloVe), Model, Output. LLMs (Large Language Models): transformer-based models trained on massive text data. Examples: GPT (OpenAI), BERT (Google), LLaMA (Meta), Gemini (Google), Claude (Anthropic). Applications: virtual assistants (Siri, Alexa), chatbots, autocomplete, spam filters. Computer Vision: enables machines to interpret visual data. Tasks: Image Classification (what is in image), Object Detection (where are objects), Image Segmentation (pixel-level), Facial Recognition, OCR. CNNs are dominant approach. Popular architectures: AlexNet, VGG, ResNet, YOLO (real-time object detection). Applications: autonomous vehicles, medical imaging, security surveillance, quality control."),
            ("AI_CH6_AI_Ethics", "IAI | CH6 - AI Ethics and Planning", "AI Ethics: responsible development and use of AI. Key issues: Bias and Fairness (AI trained on biased data produces biased outputs - facial recognition less accurate on dark skin), Privacy (AI-powered surveillance, data collection, profiling), Transparency/Explainability XAI (black box models make unexplainable decisions), Accountability (who is responsible when AI causes harm), Job Displacement (automation replacing workers), Safety and Alignment (ensuring AI behaves as intended). FATE principles: Fairness, Accountability, Transparency, Explainability. Governance: EU AI Act (2024), NIST AI Risk Management Framework. AI Planning: finding action sequences to achieve goals. STRIPS: preconditions (what must be true before action), postconditions (what becomes true/false after). Types: Classical Planning (deterministic), Probabilistic Planning (uncertain outcomes), Partial-Order Planning (flexible ordering). Reasoning under uncertainty: Bayesian Networks (probabilistic graphical models), Fuzzy Logic (degrees of truth 0-1, handles vague knowledge). Example fuzzy: temperature is 0.7 HIGH instead of just True/False. Case-Based Reasoning (CBR): solving new problems based on similar past solutions. CBR cycle: Retrieve, Reuse, Revise, Retain."),
        ]
        for name, type_, desc in core_concepts:
            self.add_concept(name, type_, desc)

    def add_concept(self, name, concept_type, description):
        self.concepts[name] = {"type": concept_type, "description": description}

    def search(self, keywords):
        results = []
        for kw in keywords:
            kw_lower = kw.lower()
            for name, info in self.concepts.items():
                if kw_lower in name.lower() or kw_lower in info.get("description", "").lower() or kw_lower in info.get("type", "").lower():
                    if not any(r['name'] == name for r in results):
                        results.append({"name": name, **info})
        return results


class BrainBotExpertSystem:
    def __init__(self):
        self.kb = SemanticNetwork()
        self.history_folder = "chat_history"
        os.makedirs(self.history_folder, exist_ok=True)
        self.conversation_memory = {}
        self.load_local_files()

    def get_history_file_path(self, user_id):
        safe_id = "".join(c if c.isalnum() else "_" for c in user_id)
        return os.path.join(self.history_folder, f"{safe_id}_history.json")

    def load_history(self, user_id):
        filepath = self.get_history_file_path(user_id)
        if os.path.exists(filepath):
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                return []
        return []

    def save_history(self, user_id, history):
        filepath = self.get_history_file_path(user_id)
        try:
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(history, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving history for {user_id}: {e}")

    def get_user_memory(self, user_id):
        if user_id not in self.conversation_memory:
            self.conversation_memory[user_id] = self.load_history(user_id)
        return self.conversation_memory[user_id]

    def add_to_history(self, user_id, role, content):
        history = self.get_user_memory(user_id)
        history.append({
            "role": role,
            "content": content,
            "timestamp": datetime.now().isoformat()
        })
        self.save_history(user_id, history)

    def extract_concepts_with_ai(self, user_input):
        prompt = (
            f"Extract the core computer science concepts from this text (fix typos automatically): '{user_input}'. "
            "Return ONLY a comma-separated list of 1 to 3 exact keywords. No chatting."
        )
        try:
            completion = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=60,
                temperature=0.1
            )
            keywords = completion.choices[0].message.content.strip().split(',')
            return [k.strip() for k in keywords if k.strip()]
        except Exception as e:
            print(f"Concept extraction error: {e}")
            return [user_input]

    def load_local_files(self):
        for filename in os.listdir(LOCAL_KB_PATH):
            path = os.path.join(LOCAL_KB_PATH, filename)
            if os.path.isfile(path):
                self.process_file(path, filename)

    def process_file(self, path, filename):
        try:
            if filename.endswith(".txt"):
                with open(path, 'r', encoding='utf-8') as f:
                    self.ingest_text(f.read(), filename)
            elif filename.endswith(".pdf"):
                reader = pypdf.PdfReader(path)
                text = "\n".join([p.extract_text() or "" for p in reader.pages])
                self.ingest_text(text, filename)
            elif filename.endswith(".pptx"):
                prs = Presentation(path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
                self.ingest_text(text, filename)
        except Exception as e:
            print(f"Error processing {filename}: {e}")

    def ingest_text(self, text, source_name):
        chunks = [text[i:i + 3000] for i in range(0, len(text), 3000)]
        for i, chunk in enumerate(chunks):
            self.kb.add_concept(f"{source_name}_chunk{i}", "document", chunk[:2800])


brainbot = BrainBotExpertSystem()


@app.route('/')
def home():
    return render_template('index.html')


@app.route('/api/chat', methods=['POST'])
def chat():
    data = request.json
    user_message = data.get('message', '').strip()
    user_id = data.get('userId', 'student_1')

    if not user_message:
        return jsonify({'reply': "Bro, what do you want to ask? 😅"}), 400

    history = brainbot.get_user_memory(user_id)
    extracted_keywords = brainbot.extract_concepts_with_ai(user_message)
    matches = brainbot.kb.search(extracted_keywords)

    kb_context = "Relevant Knowledge:\n" + "\n".join([
        f"• {m['type']}: {m['description'][:350]}" for m in matches[:5]
    ]) if matches else "No specific KB match found."

    history_text = "\n".join([
        f"{m['role'].upper()}: {m['content']}" for m in history[-7:]
    ])

    system_prompt = "You are StudBud, a friendly, slightly sarcastic APU senior who helps students ace their exams. You know ICS, ISFT, SAAD, and IAI modules inside out. Speak naturally with Malaysian uni vibe. Keep most answers short and fun but always accurate."

    user_prompt = f"""Context from Knowledge Base: {kb_context}

Recent Conversation: {history_text if history_text else "This is the start of conversation."}

Student: {user_message}

StudBud:"""

    try:
        completion = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.75,
            max_tokens=900
        )

        bot_reply = completion.choices[0].message.content.strip()

        brainbot.add_to_history(user_id, "user", user_message)
        brainbot.add_to_history(user_id, "assistant", bot_reply)

        return jsonify({'reply': bot_reply})

    except Exception as e:
        error_str = str(e).lower()
        print(f"Groq error: {e}")
        if "rate limit" in error_str or "413" in error_str or "tokens" in error_str:
            return jsonify({'reply': "Bro, I hit the token limit 😂 Give me 30 seconds then try again."}), 429
        else:
            return jsonify({'reply': "Sorry bro, something went wrong. Try again."}), 500


@app.route('/api/history', methods=['GET'])
def get_history():
    user_id = request.args.get('userId', 'student_1')
    history = brainbot.get_user_memory(user_id)
    formatted_history = [{"role": msg["role"], "content": msg["content"]} for msg in history]
    return jsonify({"history": formatted_history})


@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    save_path = os.path.join(LOCAL_KB_PATH, filename)
    file.save(save_path)
    brainbot.process_file(save_path, filename)
    return jsonify({'message': f'✅ {filename} added to StudBud knowledge base!'})


if __name__ == '__main__':
    print("StudBud Expert System is LIVE")
    print(f"Loaded {len(brainbot.kb.concepts)} concepts")
    app.run(host="0.0.0.0", port=int(os.getenv("PORT", 5000)), debug=False)
